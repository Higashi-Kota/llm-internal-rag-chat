"""Document loaders for various file formats."""

# pyright: reportUnknownMemberType=false
# pyright: reportUnknownVariableType=false
# pyright: reportUnknownArgumentType=false
# Document metadata and loader return types are not fully annotated

from collections.abc import Iterator
from pathlib import Path

from langchain_core.documents import Document

from .config import rag_settings


def load_pdf(file_path: Path) -> list[Document]:
    """Load PDF file."""
    from pypdf import PdfReader

    reader = PdfReader(str(file_path))
    documents = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text.strip():
            documents.append(
                Document(
                    page_content=text,
                    metadata={
                        "source": str(file_path),
                        "filename": file_path.name,
                        "page": i + 1,
                        "total_pages": len(reader.pages),
                    },
                )
            )
    return documents


def load_docx(file_path: Path) -> list[Document]:
    """Load DOCX file."""
    from docx import Document as DocxDocument

    doc = DocxDocument(str(file_path))
    text = "\n".join(paragraph.text for paragraph in doc.paragraphs if paragraph.text)
    if not text.strip():
        return []
    return [
        Document(
            page_content=text,
            metadata={
                "source": str(file_path),
                "filename": file_path.name,
            },
        )
    ]


def load_pptx(file_path: Path) -> list[Document]:
    """Load PPTX file."""
    from pptx import Presentation

    prs = Presentation(str(file_path))
    documents = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:  # type: ignore[attr-defined]
                texts.append(shape.text)  # type: ignore[attr-defined]
        if texts:
            documents.append(
                Document(
                    page_content="\n".join(texts),
                    metadata={
                        "source": str(file_path),
                        "filename": file_path.name,
                        "slide": i + 1,
                        "total_slides": len(prs.slides),
                    },
                )
            )
    return documents


def load_xlsx(file_path: Path) -> list[Document]:
    """Load XLSX file."""
    from openpyxl import load_workbook

    wb = load_workbook(str(file_path), read_only=True, data_only=True)
    documents = []
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        rows = []
        for row in sheet.iter_rows(values_only=True):
            row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
            if row_text.strip():
                rows.append(row_text)
        if rows:
            documents.append(
                Document(
                    page_content="\n".join(rows),
                    metadata={
                        "source": str(file_path),
                        "filename": file_path.name,
                        "sheet": sheet_name,
                    },
                )
            )
    wb.close()
    return documents


def load_txt(file_path: Path) -> list[Document]:
    """Load TXT file."""
    text = file_path.read_text(encoding="utf-8")
    if not text.strip():
        return []
    return [
        Document(
            page_content=text,
            metadata={
                "source": str(file_path),
                "filename": file_path.name,
            },
        )
    ]


# Loader mapping by extension
LoaderFunc = type[list[Document]]
LOADER_MAP: dict[str, LoaderFunc] = {
    ".pdf": load_pdf,
    ".docx": load_docx,
    ".pptx": load_pptx,
    ".xlsx": load_xlsx,
    ".txt": load_txt,
}  # type: ignore[dict-item]

SUPPORTED_EXTENSIONS = set(LOADER_MAP.keys())


def load_document(file_path: Path) -> list[Document]:
    """Load a single document based on its extension.

    Args:
        file_path: Path to the document file.

    Returns:
        List of Document objects.

    Raises:
        ValueError: If file extension is not supported.
    """
    ext = file_path.suffix.lower()
    if ext not in LOADER_MAP:
        raise ValueError(f"Unsupported file extension: {ext}")
    loader_func = LOADER_MAP[ext]
    return loader_func(file_path)  # type: ignore[operator, return-value]


def load_directory(
    directory: Path | None = None,
    recursive: bool = True,
) -> Iterator[Document]:
    """Load all documents from a directory.

    Args:
        directory: Directory path. Defaults to configured docs_dir.
        recursive: Whether to search subdirectories.

    Yields:
        Document objects from all supported files.
    """
    if directory is None:
        directory = rag_settings.docs_path

    if not directory.exists():
        raise ValueError(f"Directory does not exist: {directory}")

    pattern = "**/*" if recursive else "*"
    for file_path in directory.glob(pattern):
        if file_path.is_file() and file_path.suffix.lower() in SUPPORTED_EXTENSIONS:
            try:
                docs = load_document(file_path)
                yield from docs
            except Exception as e:
                # Log error but continue with other files
                print(f"Error loading {file_path}: {e}")
